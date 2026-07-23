from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Background #222222 (Noir Raro)
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)

# Color Tokens
RED = RGBColor(0xE6, 0x1C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AREIA = RGBColor(0xC5, 0xC0, 0xB6)
NUDE = RGBColor(0xE5, 0xE3, 0xD9)
CARD_BG = RGBColor(0x2A, 0x2A, 0x2A)

def zero_margins(tf):
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

# Header Tag (IBM Plex Mono)
tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.4))
tf = tx_box.text_frame
zero_margins(tf)
p = tf.paragraphs[0]
p.text = "LÁPIS RAЯO • BRAND INTELLIGENCE"
p.font.size = Pt(11)
p.font.color.rgb = AREIA
p.font.name = "IBM Plex Mono"
p.font.bold = True

# Header Badge (Red, IBM Plex Mono)
badge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.4), Inches(2.5), Inches(0.4))
badge_shape.fill.solid()
badge_shape.fill.fore_color.rgb = RED
badge_shape.line.fill.background()
pb = badge_shape.text_frame.paragraphs[0]
pb.text = "PREVISÃO: JAN / 2027"
pb.font.size = Pt(11)
pb.font.color.rgb = WHITE
pb.font.name = "IBM Plex Mono"
pb.font.bold = True
pb.alignment = PP_ALIGN.CENTER

# Col 1: Title & Description & Phase Badge
tx_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.3), Inches(0.9))
tf_t = tx_title.text_frame
zero_margins(tf_t)
pt = tf_t.paragraphs[0]
pt.text = "ESQUADЯO"
pt.font.size = Pt(46)
pt.font.bold = True
pt.font.color.rgb = WHITE
pt.font.name = "IBM Plex Mono"

tx_desc = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(4.3), Inches(2.2))
tf_d = tx_desc.text_frame
tf_d.word_wrap = True
zero_margins(tf_d)
pd = tf_d.paragraphs[0]
pd.text = "Centralização das especificações técnicas no Illustrator, Photoshop, After Effects e Premiere. Automação de pranchetas e composições com réguas e safe zones oficiais."
pd.font.size = Pt(14)
pd.font.color.rgb = NUDE
pd.font.name = "IBM Plex Sans"

phase_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.1), Inches(3.8), Inches(0.4))
phase_shape.fill.solid()
phase_shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
phase_shape.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
pp = phase_shape.text_frame.paragraphs[0]
pp.text = "Fase: Desenvolvimento & Checagem"
pp.font.size = Pt(11)
pp.font.color.rgb = AREIA
pp.font.name = "IBM Plex Mono"
pp.alignment = PP_ALIGN.CENTER

# Col 2: Feature Cards
features = [
    ("FLUXO GUIADO SEM DÚVIDAS", "Criação direcionada para quem não tem costume com o formato, sem consultar manuais."),
    ("ZERO REFAÇÃO POR CORTE", "Safe Zones e margens aplicadas automaticamente para proteger marcas e CTAs."),
    ("GERAÇÃO EM 1 CLIQUE", "Pranchetas individuais, múltiplas ou aplicação de guias no arquivo em andamento.")
]

y_pos = Inches(1.4)
for title, desc in features:
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), y_pos, Inches(4.4), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), y_pos, Inches(0.08), Inches(1.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()
    
    tf_f = card.text_frame
    tf_f.word_wrap = True
    zero_margins(tf_f)
    
    p1 = tf_f.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "IBM Plex Mono"
    
    p2 = tf_f.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = AREIA
    p2.font.name = "IBM Plex Sans"
    
    y_pos += Inches(1.35)

# Col 3: Demo GIF
gif_path = "/Users/gabi/Projetos_Automacoes/01-plugin-illustrator/demo.gif"
slide.shapes.add_picture(gif_path, Inches(9.8), Inches(1.4), width=Inches(2.9))

# Footer
tx_foot = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.4))
tf_ft = tx_foot.text_frame
zero_margins(tf_ft)
pft = tf_ft.paragraphs[0]
pft.text = "PROJETO 01 • AUTOMAÇÃO DE WORKFLOW                         2026 • LÁPIS RAЯO"
pft.font.size = Pt(10)
pft.font.color.rgb = AREIA
pft.font.name = "IBM Plex Mono"

pptx_path = "/Users/gabi/Projetos_Automacoes/01-plugin-illustrator/slide_esquadro.pptx"
prs.save(pptx_path)
print("Saved PPTX with exact fonts and zero-margin alignment:", pptx_path)
